"""Hybrid text-extraction pipeline — FR-OCR-01 + FR-GL-01.

Tier 0 (synchronous, seconds):
  • Docling — structure-aware PDF/DOCX/image parser that preserves tables as
    Markdown and runs its own embedded OCR engine.  Falls back gracefully when
    not installed or when parsing fails.

  • Legacy path (fallback):
      native PDF with a text layer  -> pypdf (cheap CPU path)
      scanned PDF / image           -> Tesseract OCR via PyMuPDF rasterisation
      office / text files           -> python-docx / openpyxl / python-pptx

Everything degrades gracefully: if an optional library or the Tesseract
binary is missing, the document still ingests and its status reflects why.
"""

from __future__ import annotations

import concurrent.futures
import io
import logging
import multiprocessing
import re
import tempfile
import threading
import time
import unicodedata
from dataclasses import dataclass, field
from importlib import metadata
from pathlib import Path
from typing import TypeAlias

from ..config import settings
from ..observability import emit_event
from ..utils.request_context import bound_request_context, worker_context
from . import cancellation_registry as _cancel_reg

# ── Custom exceptions ─────────────────────────────────────────────────────────


class ExtractionCancelled(RuntimeError):
    """Raised when a document is deleted mid-extraction."""


class ExtractionBatchFailed(RuntimeError):
    """Raised when a single PDF batch exhausts its per-batch retry budget."""

    def __init__(self, msg: str, last_exc: Exception | None = None) -> None:
        super().__init__(msg)
        self.last_exc = last_exc


# ── Process-pool for parallel batch extraction ────────────────────────────────
# The pool is created lazily once and kept alive so that each worker process
# warms up Docling/ONNX exactly once and then reuses the loaded models.

_extraction_pool: concurrent.futures.ProcessPoolExecutor | None = None
_pool_lock = threading.Lock()

# Subprocess-visible copy of the shared cancellation registry dict.
# Set by _pool_initializer so workers can poll without a manager round-trip.
_SHARED_REGISTRY: dict | None = None  # only populated inside subprocess workers

# ── Optional dependencies (import-guarded) ────────────────────────────────────

# Docling — layout-aware, structure-preserving parser (Tier 0 preferred path).
try:
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import (
        AcceleratorDevice,
        AcceleratorOptions,
        PdfPipelineOptions,
        RapidOcrOptions,
    )
    from docling.document_converter import (
        DocumentConverter as _DoclingConverter,
    )
    from docling.document_converter import (
        PdfFormatOption,
    )

    _HAS_DOCLING = True
except Exception:  # pragma: no cover
    _HAS_DOCLING = False

# Lazy singleton — the converter is expensive to construct (loads ~770 model
# weights). We build it once on first use and cache it for the process lifetime
# so every subsequent upload reuses the already-loaded models.
_docling_converter: "_DoclingConverter | None" = None
_docling_converter_checked: bool = False
_docling_lock = None  # initialised below after threading is available


def _rapidocr_artifacts(repo_dir: Path) -> tuple[Path, Path, Path, Path | None]:
    """Return the downloaded RapidOCR model files under ``repo_dir``."""

    onnx_files = sorted(repo_dir.rglob("*.onnx"))

    def _tokens(path: Path) -> set[str]:
        # RapidOCR's downloader currently places files directly in the model
        # root (for example ``PP-OCRv6_det_small.onnx``), while older bundles
        # used ``.../det/...`` subdirectories.  Tokenize both path components
        # and filenames so both layouts are accepted.
        return set(re.findall(r"[a-z0-9]+", path.as_posix().lower()))

    def _find(kind: str) -> Path | None:
        return next(
            (path for path in onnx_files if kind in _tokens(path)),
            None,
        )

    det = _find("det")
    cls = _find("cls")
    rec = _find("rec")
    keys = next(
        (
            path
            for path in sorted(repo_dir.rglob("*.txt"))
            if {"keys", "dict"} & _tokens(path) or "rec" in _tokens(path)
        ),
        None,
    )
    if det is None or cls is None or rec is None:
        raise FileNotFoundError(
            "RapidOCR download did not produce detector, classifier, and recognizer models"
        )
    return det, cls, rec, keys


def _ensure_rapidocr_artifacts(repo_dir: Path) -> tuple[Path, Path, Path, Path | None]:
    """Download RapidOCR models into the persistent volume when missing."""

    repo_dir.mkdir(parents=True, exist_ok=True)
    try:
        artifacts = _rapidocr_artifacts(repo_dir)
        logging.getLogger(__name__).info(
            "rapidocr: using cached artifacts under %s", repo_dir
        )
        return artifacts
    except FileNotFoundError:
        pass

    try:
        import portalocker
        import rapidocr
        import yaml
        from rapidocr import download_models
    except ImportError as exc:  # pragma: no cover - depends on runtime extras
        raise RuntimeError("RapidOCR automatic downloader is unavailable") from exc

    lock_path = repo_dir.parent / ".rapidocr-download.lock"
    with portalocker.Lock(str(lock_path), timeout=900):
        try:
            artifacts = _rapidocr_artifacts(repo_dir)
            logging.getLogger(__name__).info(
                "rapidocr: another process populated cached artifacts under %s", repo_dir
            )
            return artifacts
        except FileNotFoundError:
            pass

        package_config = Path(rapidocr.__file__).with_name("config.yaml")
        config = yaml.safe_load(package_config.read_text(encoding="utf-8"))
        config.setdefault("Global", {})["model_root_dir"] = str(repo_dir)
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".yaml", prefix="docvault-rapidocr-", delete=False
        ) as handle:
            yaml.safe_dump(config, handle, sort_keys=False)
            config_path = Path(handle.name)
        try:
            logging.getLogger(__name__).info(
                "rapidocr: downloading missing artifacts into %s", repo_dir
            )
            download_models(config_path)
        finally:
            config_path.unlink(missing_ok=True)

        return _rapidocr_artifacts(repo_dir)


def _get_docling_converter():
    """Return a process-level DocumentConverter singleton (thread-safe lazy init)."""
    global _docling_converter, _docling_converter_checked, _docling_lock
    if _docling_converter_checked:
        return _docling_converter
    import threading

    if _docling_lock is None:
        _docling_lock = threading.Lock()
    with _docling_lock:
        if _docling_converter_checked:  # double-check after acquiring lock
            return _docling_converter
        _docling_converter_checked = True
        if not _HAS_DOCLING:
            return None
        try:
            # Store all Docling model artifacts in the persistent writable volume
            # so they survive container restarts and are never downloaded into
            # the read-only /opt/venv package directory.
            _docling_artifacts_path = settings.storage_dir / "docvault-docling-artifacts"
            _docling_artifacts_path.mkdir(parents=True, exist_ok=True)

            # ── Docling layout/table model resolution ──────────────────────────
            # Docling's layout (docling-layout-heron) and table (docling-models)
            # weights are downloaded by HuggingFace into the hf_cache volume under
            # hub/models--<org>--<repo>/snapshots/<hash>/ but Docling expects them
            # under a single artifacts_path dir as <org>--<repo>/ subfolders.
            # We create a stable directory with symlinks into the HF snapshots so
            # Docling always finds its models without re-downloading them.
            import os as _os
            _hf_hub = Path(_os.environ.get("HF_HOME", "/hf_cache")) / "hub"
            _docling_models_dir = settings.storage_dir / "docling-models"
            _docling_models_dir.mkdir(parents=True, exist_ok=True)
            _LAYOUT_REPO = "docling-project--docling-layout-heron"
            _MODELS_REPO = "docling-project--docling-models"
            for _repo_name, _hf_prefix in (
                (_LAYOUT_REPO, "models--docling-project--docling-layout-heron"),
                (_MODELS_REPO, "models--docling-project--docling-models"),
            ):
                _link = _docling_models_dir / _repo_name
                if _link.is_symlink():
                    # Re-evaluate: if the target has no weights, re-link to best snap
                    _current_target = Path(_os.readlink(str(_link)))
                    _has_weights = bool(list(_current_target.rglob("*.safetensors")))
                    if _has_weights:
                        continue  # already pointing to a good snapshot
                    _os.unlink(str(_link))
                    logging.getLogger(__name__).warning(
                        "docling: re-linking %s (no weights in current target)", _repo_name
                    )
                if not _link.exists():
                    _hf_model_dir = _hf_hub / _hf_prefix
                    _snaps_dir = _hf_model_dir / "snapshots"
                    # Pick the snapshot with the most .safetensors weight files
                    _best_snap = None
                    _best_count = -1
                    if _snaps_dir.is_dir():
                        for _snap in _snaps_dir.iterdir():
                            _count = len(list(_snap.rglob("*.safetensors")))
                            if _count > _best_count:
                                _best_count = _count
                                _best_snap = _snap
                    if _best_snap:
                        _os.symlink(str(_best_snap), str(_link))
                        logging.getLogger(__name__).info(
                            "docling: linked %s -> %s (%d weight files)",
                            _link.name, _best_snap, _best_count,
                        )

            # Use the onnxruntime backend — it is installed via docling[rapidocr]
            # in the ai extra and works on every deployment, including those
            # without the visual extra (which is the only place torch lives).
            #
            _repo_dir = _docling_artifacts_path / "RapidOcr"
            _det, _cls, _rec, _keys = _ensure_rapidocr_artifacts(_repo_dir)
            _font_path = Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
            _ocr_opts = RapidOcrOptions(
                lang=["english"],
                backend="onnxruntime",
                det_model_path=str(_det),
                cls_model_path=str(_cls),
                rec_model_path=str(_rec),
                rec_keys_path=str(_keys) if _keys is not None else None,
                font_path=str(_font_path) if _font_path.is_file() else None,
            )

            # Pass the explicit artifacts_path so Docling always uses the locally
            # cached layout/table models and never falls back to downloading.
            _resolved_artifacts = (
                _docling_models_dir
                if (_docling_models_dir / _LAYOUT_REPO).exists()
                else None
            )
            _pdf_pipeline_options = PdfPipelineOptions(
                do_ocr=True,
                do_table_structure=True,
                artifacts_path=_resolved_artifacts,
                ocr_options=_ocr_opts,
                accelerator_options=AcceleratorOptions(
                    device=AcceleratorDevice.CPU,
                ),
            )
            _allowed_formats = [
                InputFormat.PDF,
                InputFormat.IMAGE,
                InputFormat.DOCX,
                InputFormat.PPTX,
                InputFormat.XLSX,
                InputFormat.HTML,
                InputFormat.MD,
                InputFormat.CSV,
            ]
            for _candidate in ("ASCIIDOC", "EPUB", "ODT", "ODS", "ODP", "DOC", "XLS", "PPT", "LATEX", "VTT"):
                if hasattr(InputFormat, _candidate):
                    _allowed_formats.append(getattr(InputFormat, _candidate))

            _has_whisper = False
            try:
                import whisper
                _has_whisper = True
            except Exception:
                pass

            _format_options = {
                InputFormat.PDF: PdfFormatOption(
                    pipeline_options=_pdf_pipeline_options
                )
            }
            if _has_whisper:
                try:
                    from docling.document_converter import AudioFormatOption, VideoFormatOption
                    from docling.datamodel.pipeline_options import AsrPipelineOptions, VideoPipelineOptions
                    if hasattr(InputFormat, "AUDIO"):
                        _allowed_formats.append(getattr(InputFormat, "AUDIO"))
                        _format_options[getattr(InputFormat, "AUDIO")] = AudioFormatOption(
                            pipeline_options=AsrPipelineOptions()
                        )
                    if hasattr(InputFormat, "VIDEO"):
                        _allowed_formats.append(getattr(InputFormat, "VIDEO"))
                        _format_options[getattr(InputFormat, "VIDEO")] = VideoFormatOption(
                            pipeline_options=VideoPipelineOptions()
                        )
                except Exception:
                    pass

            _docling_converter = _DoclingConverter(
                allowed_formats=_allowed_formats,
                format_options=_format_options,
            )

            emit_event(
                "extraction.engine.initialized",
                level=logging.INFO,
                component="extraction",
                operation="docling",
            )
        except Exception as exc:
            logging.error(f"Docling initialization failed: {exc}", exc_info=True)
            emit_event(
                "extraction.engine.initialized",
                level=logging.ERROR,
                component="extraction",
                operation="docling",
                outcome="error",
                error=exc,
            )
            _docling_converter = None
    return _docling_converter


# Legacy Tesseract stack
try:
    import pytesseract
    from PIL import Image

    _HAS_TESSERACT_LIB = True
except Exception:  # pragma: no cover
    _HAS_TESSERACT_LIB = False

try:
    import fitz  # PyMuPDF

    _HAS_FITZ = True
except Exception:  # pragma: no cover
    _HAS_FITZ = False

try:
    from pypdf import PdfReader

    _HAS_PYPDF = True
except Exception:  # pragma: no cover
    _HAS_PYPDF = False

AUDIO_EXTS = {".wav", ".mp3", ".m4a", ".aac", ".ogg", ".flac"}
VIDEO_EXTS = {".mp4", ".avi", ".mov", ".mkv", ".webm"}
MEDIA_EXTS = AUDIO_EXTS | VIDEO_EXTS
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp"}
TEXT_EXTS = {
    ".txt",
    ".md",
    ".csv",
    ".log",
    ".json",
    ".html",
    ".htm",
    ".xhtml",
    ".xml",
    ".yaml",
    ".yml",
    ".vtt",
    ".asciidoc",
    ".adoc",
    ".tex",
    ".latex",
}
DOCLING_EXTS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".html",
    ".htm",
    ".xhtml",
    ".md",
    ".csv",
    ".png",
    ".jpg",
    ".jpeg",
    ".tif",
    ".tiff",
    ".bmp",
    ".webp",
    ".gif",
    ".wav",
    ".mp3",
    ".m4a",
    ".aac",
    ".ogg",
    ".flac",
    ".mp4",
    ".avi",
    ".mov",
    ".mkv",
    ".webm",
    ".vtt",
    ".epub",
    ".asciidoc",
    ".adoc",
}
# Threshold of extractable chars per page below which a PDF is deemed "scanned".
_NATIVE_TEXT_MIN_CHARS = 40
EXTRACTION_PIPELINE_VERSION = "hybrid-v2"

QualitySignal: TypeAlias = bool | float | int | str


class _OcrLanguageUnavailable(RuntimeError):
    """Configured OCR language packs are absent from the worker runtime."""


def _component_version(distribution: str, fallback: str) -> str:
    try:
        value = metadata.version(distribution)
    except metadata.PackageNotFoundError:
        value = fallback
    normalized = re.sub(r"[^A-Za-z0-9_.+-]", "-", value).strip("-")
    return (normalized or fallback)[:40]


def _detect_language(text: str) -> str:
    latin = sum(
        1
        for character in text
        if ("LATIN" in unicodedata.name(character, "") and character.isalpha())
    )
    devanagari = sum(1 for character in text if "\u0900" <= character <= "\u097f")
    if latin and devanagari:
        return "eng+hin"
    if devanagari:
        return "hin"
    if latin:
        return "eng"
    return "und"


def _measure_quality(
    text: str,
    *,
    page_count: int,
    ocr_confidence: float | None,
) -> tuple[float, dict[str, QualitySignal]]:
    character_count = len(text)
    non_whitespace = sum(1 for character in text if not character.isspace())
    printable = sum(1 for character in text if character.isprintable())
    replacements = text.count("\ufffd")
    pages = max(page_count, 1)
    printable_ratio = printable / character_count if character_count else 0.0
    replacement_ratio = replacements / character_count if character_count else 0.0
    characters_per_page = non_whitespace / pages
    density_score = min(1.0, characters_per_page / 80.0)
    cleanliness_score = max(0.0, printable_ratio - (replacement_ratio * 4.0))
    if ocr_confidence is None:
        score = (0.6 * cleanliness_score) + (0.4 * density_score)
    else:
        score = (
            (0.6 * max(0.0, min(1.0, ocr_confidence)))
            + (0.25 * cleanliness_score)
            + (0.15 * density_score)
        )
    signals: dict[str, QualitySignal] = {
        "character_count": character_count,
        "non_whitespace_characters": non_whitespace,
        "characters_per_page": round(characters_per_page, 3),
        "printable_ratio": round(printable_ratio, 4),
        "replacement_character_ratio": round(replacement_ratio, 4),
        "ocr_confidence_measured": ocr_confidence is not None,
    }
    return round(max(0.0, min(1.0, score)), 3), signals


@dataclass
class OcrResult:
    text: str = ""
    status: str = "pending"  # native | ocr | unavailable | skipped | error
    confidence: float | None = None
    page_count: int = 0
    language: str = "und"
    notes: list[str] = field(default_factory=list)
    extractor_name: str = "none"
    extractor_version: str = EXTRACTION_PIPELINE_VERSION
    ocr_engine: str | None = None
    ocr_engine_version: str | None = None
    ocr_languages: list[str] = field(default_factory=list)
    missing_ocr_languages: list[str] = field(default_factory=list)
    quality_score: float | None = None
    quality_signals: dict[str, QualitySignal] = field(default_factory=dict)


def _finalize_result(result: OcrResult) -> OcrResult:
    result.text = result.text.replace("\x00", "")
    result.language = _detect_language(result.text)
    result.quality_score, measured = _measure_quality(
        result.text,
        page_count=result.page_count,
        ocr_confidence=result.confidence,
    )
    result.quality_signals = {
        **measured,
        **result.quality_signals,
        "language": result.language,
        "page_count": max(0, result.page_count),
        "requested_ocr_languages": "+".join(settings.ocr_languages),
        "effective_ocr_languages": "+".join(result.ocr_languages),
        "missing_ocr_languages": "+".join(result.missing_ocr_languages),
        "ocr_language_pack_complete": not result.missing_ocr_languages,
    }
    return result


# ── Docling (Tier 0) ──────────────────────────────────────────────────────────


def _docling_available() -> bool:
    return _HAS_DOCLING and settings.use_docling


def _extract_with_docling(path: Path) -> OcrResult:
    """Use the cached Docling converter to produce structure-aware Markdown."""
    res = OcrResult(
        extractor_name="docling",
        extractor_version=_component_version("docling", EXTRACTION_PIPELINE_VERSION),
    )
    converter = _get_docling_converter()
    if converter is None:
        res.status = "error"
        res.notes.append("Docling converter not available")
        return res
    try:
        doc = converter.convert(str(path))
        # Export as Markdown — tables become pipe-tables, headings are preserved.
        md = doc.document.export_to_markdown()
        res.text = md or ""
        _num_pages = getattr(doc.document, "num_pages", 0)
        res.page_count = (
            int(_num_pages()) if callable(_num_pages) else int(_num_pages or 0)
        ) or 1
        # Docling does not expose one calibrated document-level confidence.
        # Parser success is therefore recorded through provenance and measured
        # text-quality signals, never as a fabricated confidence of 1.0.
        res.status = "native"
        res.confidence = None
        res.notes.append("parsed by Docling (layout-aware)")
    except Exception as exc:  # pragma: no cover
        res.status = "error"
        res.notes.append(f"Docling failed: {exc}")
    return res


# ── Legacy Tesseract stack ────────────────────────────────────────────────────


def _tesseract_available() -> bool:
    if not _HAS_TESSERACT_LIB:
        return False
    try:
        pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


def _effective_tesseract_languages(
    available: set[str] | None = None,
) -> tuple[list[str], list[str]]:
    requested = list(settings.ocr_languages)
    if available is None:
        try:
            available = set(pytesseract.get_languages(config=""))
        except Exception:
            available = set()
    effective = [language for language in requested if language in available]
    missing = [language for language in requested if language not in available]
    return effective, missing


def _ocr_image_bytes(
    data: bytes,
) -> tuple[str, float | None, list[str], list[str]]:
    """Run Tesseract and report measured confidence and language provenance."""
    image = Image.open(io.BytesIO(data))
    languages, missing_languages = _effective_tesseract_languages()
    if not languages:
        raise _OcrLanguageUnavailable
    language_argument = "+".join(languages)
    text = pytesseract.image_to_string(image, lang=language_argument)
    conf = None
    try:
        d = pytesseract.image_to_data(
            image,
            lang=language_argument,
            output_type=pytesseract.Output.DICT,
        )
        vals = [
            int(c)
            for c in d.get("conf", [])
            if str(c).lstrip("-").isdigit() and int(c) >= 0
        ]
        if vals:
            conf = round(sum(vals) / len(vals) / 100.0, 3)
    except Exception:
        pass
    return text, conf, languages, missing_languages


def _extract_pdf_legacy(path: Path) -> OcrResult:
    res = OcrResult(
        extractor_name="pypdf",
        extractor_version=_component_version("pypdf", EXTRACTION_PIPELINE_VERSION),
    )

    # 1) Try native text layer first (cheap CPU path).
    native_text = ""
    pages = 0
    if _HAS_PYPDF:
        try:
            reader = PdfReader(str(path))
            pages = len(reader.pages)
            native_text = "\n".join((p.extract_text() or "") for p in reader.pages)
        except Exception as exc:
            res.notes.append(f"pypdf failed: {exc}")

    res.page_count = pages
    # Heuristic: enough extractable text (scaled by page count) => digital PDF.
    if native_text.strip() and len(native_text.strip()) >= _NATIVE_TEXT_MIN_CHARS * max(
        pages, 1
    ):
        res.text = native_text
        res.status = "native"
        res.confidence = None
        return res

    # 2) Scanned PDF -> rasterize pages and OCR them.
    if not _tesseract_available():
        res.text = native_text
        res.status = "unavailable" if not native_text.strip() else "native"
        res.notes.append("Tesseract binary not available for scanned pages")
        return res
    if not _HAS_FITZ:
        res.text = native_text
        res.status = "unavailable" if not native_text.strip() else "native"
        res.notes.append("PyMuPDF not installed; cannot rasterize scanned PDF")
        return res

    try:
        doc = fitz.open(str(path))
        res.page_count = doc.page_count
        chunks, confs = [], []
        effective_languages: set[str] = set()
        missing_languages: set[str] = set()
        for page in doc:
            pix = page.get_pixmap(dpi=200)
            text, conf, languages, missing = _ocr_image_bytes(pix.tobytes("png"))
            chunks.append(text)
            if conf is not None:
                confs.append(conf)
            effective_languages.update(languages)
            missing_languages.update(missing)
        res.text = "\n".join(chunks)
        res.status = "ocr"
        res.confidence = round(sum(confs) / len(confs), 3) if confs else None
        res.extractor_name = "pymupdf+tesseract"
        res.extractor_version = _component_version(
            "PyMuPDF",
            EXTRACTION_PIPELINE_VERSION,
        )
        res.ocr_engine = "tesseract"
        res.ocr_engine_version = str(pytesseract.get_tesseract_version())[:40]
        res.ocr_languages = sorted(effective_languages)
        res.missing_ocr_languages = sorted(missing_languages)
        if res.missing_ocr_languages:
            res.notes.append("one or more configured OCR language packs are unavailable")
    except _OcrLanguageUnavailable:
        res.text = native_text
        res.status = "unavailable"
        res.missing_ocr_languages = list(settings.ocr_languages)
        res.notes.append("configured OCR language packs are unavailable")
    except Exception as exc:
        res.text = native_text
        res.status = "error"
        res.notes.append(f"OCR failed: {exc}")
    return res


def _extract_image_legacy(path: Path) -> OcrResult:
    res = OcrResult(
        page_count=1,
        extractor_name="pillow+tesseract",
        extractor_version=_component_version("Pillow", EXTRACTION_PIPELINE_VERSION),
    )
    if not _tesseract_available():
        res.status = "unavailable"
        res.notes.append("Tesseract binary not available")
        return res
    try:
        text, conf, languages, missing = _ocr_image_bytes(path.read_bytes())
        res.text, res.confidence, res.status = text, conf, "ocr"
        res.ocr_engine = "tesseract"
        res.ocr_engine_version = str(pytesseract.get_tesseract_version())[:40]
        res.ocr_languages = languages
        res.missing_ocr_languages = missing
        if missing:
            res.notes.append("one or more configured OCR language packs are unavailable")
    except _OcrLanguageUnavailable:
        res.status = "unavailable"
        res.missing_ocr_languages = list(settings.ocr_languages)
        res.notes.append("configured OCR language packs are unavailable")
    except Exception as exc:
        res.status = "error"
        res.notes.append(f"OCR failed: {exc}")
    return res


def _extract_office_or_text(path: Path, extension: str | None = None) -> OcrResult:
    res = OcrResult(page_count=1, confidence=None)
    ext = (extension or path.suffix).lower()
    try:
        if ext in TEXT_EXTS:
            res.text = path.read_text(errors="ignore")
            res.extractor_name = "plain-text"
            res.extractor_version = EXTRACTION_PIPELINE_VERSION
        elif ext == ".docx":
            import docx  # python-docx

            res.extractor_name = "python-docx"
            res.extractor_version = _component_version(
                "python-docx",
                EXTRACTION_PIPELINE_VERSION,
            )
            doc = docx.Document(str(path))
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for table in doc.tables:
                # Render tables as Markdown for better structure preservation.
                table_rows = [
                    [cell.text.strip() for cell in row.cells] for row in table.rows
                ]
                if table_rows:
                    header = "| " + " | ".join(table_rows[0]) + " |"
                    separator = "| " + " | ".join(["---"] * len(table_rows[0])) + " |"
                    body = "\n".join(
                        "| " + " | ".join(row) + " |" for row in table_rows[1:]
                    )
                    parts.append("\n".join([header, separator, body]))
            res.text = "\n\n".join(parts)
        elif ext == ".xlsx":
            import openpyxl

            res.extractor_name = "openpyxl"
            res.extractor_version = _component_version(
                "openpyxl",
                EXTRACTION_PIPELINE_VERSION,
            )
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            sheet_lines: list[str] = []
            for ws in wb.worksheets:
                sheet_lines.append(f"## Sheet: {ws.title}")
                sheet_rows = list(ws.iter_rows(values_only=True))
                if sheet_rows:
                    # Render as Markdown table.
                    header = sheet_rows[0]
                    sheet_lines.append(
                        "| "
                        + " | ".join(
                            str(cell) if cell is not None else "" for cell in header
                        )
                        + " |"
                    )
                    sheet_lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                    for row in sheet_rows[1:]:
                        sheet_lines.append(
                            "| "
                            + " | ".join(
                                "" if cell is None else str(cell) for cell in row
                            )
                            + " |"
                        )
            res.text = "\n".join(sheet_lines)
        elif ext == ".pptx":
            from pptx import Presentation

            res.extractor_name = "python-pptx"
            res.extractor_version = _component_version(
                "python-pptx",
                EXTRACTION_PIPELINE_VERSION,
            )
            prs = Presentation(str(path))
            chunks, slides = [], 0
            for slide in prs.slides:
                slides += 1
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
                    if shape.has_table:
                        for r in shape.table.rows:
                            chunks.append(" | ".join(c.text for c in r.cells))
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    chunks.append(slide.notes_slide.notes_text_frame.text)
            res.text = "\n".join(c for c in chunks if c and c.strip())
            res.page_count = slides
        else:
            res.status = "skipped"
            res.notes.append(
                f"Stored, but no text extractor for {ext} (still searchable by title)"
            )
            return res
        res.status = "native"
    except Exception as exc:
        res.status = "error"
        res.notes.append(f"text extraction failed: {exc}")
    return res


# ── Parallel batch extraction ─────────────────────────────────────────────────


def _pool_initializer(shared_registry: dict) -> None:
    """Run once in each new pool worker process to cache the registry ref."""
    global _SHARED_REGISTRY
    _SHARED_REGISTRY = shared_registry


def _get_extraction_pool() -> concurrent.futures.ProcessPoolExecutor:
    """Return the lazily-created, process-reusing extraction pool."""
    global _extraction_pool
    if _extraction_pool is not None:
        return _extraction_pool
    with _pool_lock:
        if _extraction_pool is not None:
            return _extraction_pool
        # forkserver avoids inheriting file-descriptor leaks and lock state
        # that would occur with the default "fork" start method on Linux.
        ctx = multiprocessing.get_context("forkserver")
        _extraction_pool = concurrent.futures.ProcessPoolExecutor(
            max_workers=settings.extraction_max_workers,
            mp_context=ctx,
            initializer=_pool_initializer,
            initargs=(_cancel_reg.get_registry(),),
        )
        logging.getLogger(__name__).info(
            "extraction_pool: started with max_workers=%d",
            settings.extraction_max_workers,
        )
    return _extraction_pool


def _reset_extraction_pool() -> None:
    """Tear down a broken pool so the next call to _get_extraction_pool rebuilds it."""
    global _extraction_pool
    with _pool_lock:
        old = _extraction_pool
        _extraction_pool = None
    if old is not None:
        try:
            for p in getattr(old, "_processes", {}).values():
                try:
                    p.terminate()
                except Exception:
                    pass
            old.shutdown(wait=False, cancel_futures=True)
        except Exception:
            pass
        logging.getLogger(__name__).warning(
            "extraction_pool: torn down (will respawn on next batch submission)"
        )


def _split_pdf_into_batches(path: Path, batch_size: int) -> list[Path]:
    """Split *path* into ≤batch_size-page PDFs stored in a temp directory.

    Returns a list of Paths.  The caller is responsible for cleaning up those
    files when they are no longer needed.
    """
    import os as _os

    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(str(path))
    total = len(reader.pages)
    batch_dir = settings.storage_dir / ".pdf-batches"
    batch_dir.mkdir(parents=True, exist_ok=True)
    stem = path.stem[:20]
    pid = _os.getpid()
    paths: list[Path] = []
    try:
        for start in range(0, total, batch_size):
            writer = PdfWriter()
            for i in range(start, min(start + batch_size, total)):
                writer.add_page(reader.pages[i])
            batch_path = batch_dir / f"batch_{pid}_{start}_{stem}.pdf"
            with batch_path.open("wb") as fh:
                writer.write(fh)
            paths.append(batch_path)
    except Exception:
        # Clean up anything that was already written.
        for p in paths:
            p.unlink(missing_ok=True)
        raise
    return paths


# Top-level (picklable) batch-worker function.  Must NOT be a closure or method.
def _extract_batch_worker(
    batch_path_str: str,
    document_id: int,
    filename: str,
) -> dict:
    """Run inside a pool subprocess: extract one batch PDF and return a dict.

    Returns a plain dict (not OcrResult) so it can be pickled across the
    process boundary without importing the full module graph in each worker.
    """
    if _SHARED_REGISTRY is not None and _SHARED_REGISTRY.get(document_id, False):
        return {"status": "cancelled", "text": "", "page_count": 0, "notes": [], "confidence": None}

    path = Path(batch_path_str)
    if not path.exists():
        return {
            "status": "error",
            "text": "",
            "page_count": 0,
            "confidence": None,
            "notes": [f"Batch file {path.name} not found on disk"],
            "extractor_name": "docling",
            "extractor_version": EXTRACTION_PIPELINE_VERSION,
        }
    result = _extract_with_docling(path)
    return {
        "status": result.status,
        "text": result.text,
        "page_count": result.page_count,
        "confidence": result.confidence,
        "notes": result.notes,
        "extractor_name": result.extractor_name,
        "extractor_version": result.extractor_version,
    }


def _dict_to_ocr_result(d: dict) -> OcrResult:
    """Reconstruct an OcrResult from the dict returned by _extract_batch_worker."""
    r = OcrResult()
    r.status = d.get("status", "error")
    r.text = d.get("text", "")
    r.page_count = d.get("page_count", 0)
    r.confidence = d.get("confidence")
    r.notes = list(d.get("notes", []))
    r.extractor_name = d.get("extractor_name", "docling-batch")
    r.extractor_version = d.get("extractor_version", EXTRACTION_PIPELINE_VERSION)
    return r


def _backoff(attempt: int, max_seconds: float = 30.0) -> None:
    time.sleep(min(2 ** attempt, max_seconds))


def _extract_batch_with_retry(
    pool: concurrent.futures.ProcessPoolExecutor,
    batch_path: Path,
    document_id: int,
    filename: str,
    max_retries: int,
) -> OcrResult:
    """Submit one batch to the pool with per-batch retry on failure.

    Retries only the *failing* batch; all other batches are unaffected.
    Raises ``ExtractionCancelled`` if the document is deleted mid-retry.
    Raises ``ExtractionBatchFailed`` after *max_retries* exhausted.
    """
    log = logging.getLogger(__name__)
    last_exc: Exception | None = None

    for attempt in range(1, max_retries + 1):
        # Honour cancellation before every attempt.
        if _cancel_reg.is_cancelled(document_id):
            raise ExtractionCancelled(f"document {document_id} cancelled before batch retry {attempt}")

        try:
            future = pool.submit(
                _extract_batch_worker,
                str(batch_path),
                document_id,
                filename,
            )
            result_dict = future.result(timeout=settings.max_extraction_seconds)

        except concurrent.futures.BrokenProcessPool as exc:
            log.warning(
                "extraction_pool: BrokenProcessPool on attempt %d/%d for %s — resetting pool",
                attempt, max_retries, batch_path.name,
            )
            _reset_extraction_pool()
            # Rebuild pool and get fresh reference for subsequent attempts.
            pool = _get_extraction_pool()
            last_exc = exc
            _backoff(attempt)
            continue

        except (TimeoutError, OSError, Exception) as exc:
            log.warning(
                "Batch %s attempt %d/%d failed: %s",
                batch_path.name, attempt, max_retries, exc,
            )
            last_exc = exc
            _backoff(attempt)
            continue

        # Worker reported cancellation.
        if result_dict.get("status") == "cancelled":
            raise ExtractionCancelled(f"document {document_id} cancelled inside worker")

        # Docling soft-failure — retry if budget remains.
        if result_dict.get("status") == "error" and attempt < max_retries:
            log.warning(
                "Batch %s Docling error (attempt %d/%d) — requeueing",
                batch_path.name, attempt, max_retries,
            )
            _backoff(attempt)
            continue

        # Success (or final attempt — return whatever we got).
        return _dict_to_ocr_result(result_dict)

    raise ExtractionBatchFailed(
        f"Batch {batch_path.name} failed after {max_retries} attempts",
        last_exc=last_exc,
    )


def _merge_batch_results(results: list[OcrResult]) -> OcrResult:
    """Combine per-batch OcrResults into one unified result."""
    if not results:
        return OcrResult(status="error", notes=["no batch results to merge"])

    merged = OcrResult(
        extractor_name="docling-parallel",
        extractor_version=results[0].extractor_version,
    )
    texts: list[str] = []
    notes: list[str] = []
    total_pages = 0
    conf_sum = 0.0
    conf_pages = 0

    for r in results:
        texts.append(r.text)
        notes.extend(r.notes)
        total_pages += r.page_count
        if r.confidence is not None and r.page_count > 0:
            conf_sum += r.confidence * r.page_count
            conf_pages += r.page_count

    merged.text = "\n\n".join(t for t in texts if t)
    merged.page_count = total_pages
    merged.confidence = round(conf_sum / conf_pages, 3) if conf_pages else None
    merged.notes = notes
    merged.status = "native"
    return merged


def _extract_with_docling_parallel(path: Path, document_id: int) -> OcrResult:
    """Split a large PDF into batches and extract them in parallel subprocesses.

    Each batch is individually retried on failure.  All batches that succeed
    are merged into a single OcrResult.  If any batch exhausts its retry budget
    the whole extraction fails with ExtractionBatchFailed.  If the document is
    deleted the whole extraction is cancelled with ExtractionCancelled.
    """
    log = logging.getLogger(__name__)
    batch_size = settings.pdf_batch_size
    max_retries = settings.pdf_batch_max_retries
    filename = path.name

    batch_paths: list[Path] = []
    try:
        batch_paths = _split_pdf_into_batches(path, batch_size)
        log.info(
            "extraction: document %d split into %d batch(es) of ≤%d pages",
            document_id, len(batch_paths), batch_size,
        )

        if len(batch_paths) == 1:
            # Single batch — skip pool overhead and call Docling directly.
            result = _extract_with_docling(batch_paths[0])
            return result

        pool = _get_extraction_pool()

        # Submit all batches concurrently to the pool.  Retry logic runs in
        # *this* process (not inside a subprocess) so pool references and
        # non-picklable objects can be used freely.
        batch_futures: dict[concurrent.futures.Future, tuple[Path, int]] = {}

        def _submit_batch(bp: Path) -> concurrent.futures.Future:
            return pool.submit(_extract_batch_worker, str(bp), document_id, filename)

        for bp in batch_paths:
            if _cancel_reg.is_cancelled(document_id):
                raise ExtractionCancelled(f"document {document_id} cancelled before submit")
            batch_futures[_submit_batch(bp)] = (bp, 0)

        ordered_results: dict[Path, OcrResult] = {}

        # Timeout per batch step: if no batch completes within this window, the pool is considered stalled.
        batch_wait_timeout = max(300.0, float(settings.max_extraction_seconds))
        # Total safety deadline: scale with batch count (at least 90s per batch, minimum 30 minutes).
        total_deadline = time.monotonic() + max(1800.0, len(batch_paths) * 90.0)

        while batch_futures:
            if time.monotonic() > total_deadline:
                for f in batch_futures:
                    f.cancel()
                _reset_extraction_pool()
                raise ExtractionBatchFailed(
                    f"Document {document_id}: overall extraction time exceeded safety limit ({len(batch_paths)} batches)"
                )

            # Wait for any future to complete.
            done, _ = concurrent.futures.wait(
                list(batch_futures.keys()),
                timeout=batch_wait_timeout,
                return_when=concurrent.futures.FIRST_COMPLETED,
            )

            if _cancel_reg.is_cancelled(document_id):
                for f in batch_futures:
                    f.cancel()
                _reset_extraction_pool()
                raise ExtractionCancelled(f"document {document_id} cancelled while awaiting batches")

            if not done:
                # Progress stall timeout — cancel everything and tear down worker subprocesses.
                for f in batch_futures:
                    f.cancel()
                _reset_extraction_pool()
                raise ExtractionBatchFailed(
                    f"Document {document_id}: batch pool stalled (no batch completed within {int(batch_wait_timeout)}s)"
                )

            for fut in done:
                bp, attempt = batch_futures.pop(fut)
                try:
                    result_dict = fut.result()
                except concurrent.futures.BrokenProcessPool as exc:
                    log.warning("extraction_pool: BrokenProcessPool — resetting and requeueing %s", bp.name)
                    _reset_extraction_pool()
                    pool = _get_extraction_pool()
                    attempt += 1
                    if attempt >= max_retries:
                        raise ExtractionBatchFailed(
                            f"Batch {bp.name} failed after {max_retries} attempts (BrokenProcessPool)",
                            last_exc=exc,
                        )
                    _backoff(attempt)
                    batch_futures[_submit_batch(bp)] = (bp, attempt)
                    continue
                except Exception as exc:
                    attempt += 1
                    if attempt >= max_retries:
                        raise ExtractionBatchFailed(
                            f"Batch {bp.name} failed after {max_retries} attempts",
                            last_exc=exc,
                        )
                    log.warning("Batch %s attempt %d/%d failed: %s — requeueing", bp.name, attempt, max_retries, exc)
                    _backoff(attempt)
                    batch_futures[_submit_batch(bp)] = (bp, attempt)
                    continue

                if result_dict.get("status") == "cancelled":
                    raise ExtractionCancelled(f"document {document_id} cancelled inside worker")

                if result_dict.get("status") == "error" and attempt < max_retries - 1:
                    attempt += 1
                    log.warning("Batch %s Docling error attempt %d/%d — requeueing", bp.name, attempt, max_retries)
                    _backoff(attempt)
                    batch_futures[_submit_batch(bp)] = (bp, attempt)
                    continue

                log.info("extraction: document %d batch %s complete", document_id, bp.name)
                ordered_results[bp] = _dict_to_ocr_result(result_dict)

        # Merge in original split order.
        results_in_order = [ordered_results[bp] for bp in batch_paths if bp in ordered_results]
        merged = _merge_batch_results(results_in_order)
        log.info(
            "extraction: document %d merged %d batch(es) → %d chars, %d pages",
            document_id, len(results_in_order), len(merged.text), merged.page_count,
        )
        return merged

    finally:
        # If any batch futures remain unfinished (e.g. timeout, cancellation, or error),
        # terminate the pool processes first so orphan background workers cannot access unlinked files.
        if batch_futures:
            for f in batch_futures:
                f.cancel()
            _reset_extraction_pool()

        # Always clean up temp batch files.
        for bp in batch_paths:
            try:
                bp.unlink(missing_ok=True)
            except Exception:
                pass


def _extract_vtt(path: Path, filename: str = "") -> OcrResult:
    """Extract text and timestamp breakdown from WebVTT (.vtt) subtitle / caption file."""
    display_name = filename or path.name
    res = OcrResult(
        extractor_name="webvtt-parser",
        extractor_version=EXTRACTION_PIPELINE_VERSION,
        page_count=1,
        quality_score=1.0,
        quality_signals={"character_count": 0, "ocr_confidence_measured": False},
    )
    try:
        content = path.read_text(encoding="utf-8", errors="replace")
        lines = content.splitlines()
        cue_lines = []
        current_time = ""
        current_text = []

        for line in lines:
            trimmed = line.strip()
            if not trimmed or trimmed.startswith("WEBVTT") or trimmed.startswith("NOTE"):
                continue
            if "-->" in trimmed:
                if current_time and current_text:
                    cue_lines.append(f"- `[{current_time}]` {' '.join(current_text)}")
                    current_text = []
                parts = trimmed.split("-->")
                start_part = parts[0].strip().split(".")[0]
                end_part = parts[1].strip().split()[0].split(".")[0]
                current_time = f"{start_part} - {end_part}"
                continue
            if re.match(r"^\d+$", trimmed):
                continue
            cleaned = re.sub(r"<[^>]+>", "", trimmed)
            if cleaned:
                current_text.append(cleaned)

        if current_time and current_text:
            cue_lines.append(f"- `[{current_time}]` {' '.join(current_text)}")

        if cue_lines:
            res.text = f"# Video Subtitles: {display_name}\n\n## Timestamps & Dialogue\n" + "\n".join(cue_lines)
        else:
            plain_text = [
                re.sub(r"<[^>]+>", "", l.strip())
                for l in lines
                if l.strip() and not l.startswith("WEBVTT") and "-->" not in l
            ]
            res.text = "\n".join(plain_text)

        res.status = "native"
        res.notes.append("parsed WebVTT subtitles with timestamp breakdown")
    except Exception as exc:
        res.status = "error"
        res.notes.append(f"WebVTT extraction failed: {exc}")
    return res


def _extract_media_fallback(path: Path, ext: str, filename: str = "") -> OcrResult:
    """Extract or transcribe audio/video using Whisper with timestamp breakdown or ffprobe metadata fallback."""
    display_name = filename or path.name
    res = OcrResult(
        extractor_name="media-asr",
        extractor_version=EXTRACTION_PIPELINE_VERSION,
        page_count=1,
        quality_score=1.0,
        quality_signals={"character_count": 0, "ocr_confidence_measured": False},
    )

    # 1. Try whisper ASR transcription with timestamps if installed
    try:
        import os
        import shutil
        if "/data/bin" not in os.environ.get("PATH", ""):
            os.environ["PATH"] = f"/data/bin:{os.environ.get('PATH', '')}"
        if not shutil.which("ffmpeg"):
            try:
                import imageio_ffmpeg
                ffmpeg_exe = imageio_ffmpeg.get_ffmpeg_exe()
                ffmpeg_dir = str(Path(ffmpeg_exe).parent)
                os.environ["PATH"] = f"{ffmpeg_dir}:/data/bin:{os.environ.get('PATH', '')}"
            except Exception:
                pass

        import whisper

        model = whisper.load_model("tiny")
        transcription = model.transcribe(str(path))
        segments = transcription.get("segments", [])
        if segments:
            lines = [f"# Media Transcript: {display_name}"]
            lang = transcription.get("language", "und")
            if lang != "und":
                lines.append(f"- **Language**: {lang}")
            lines.append("\n## Timestamps & Dialogue\n")
            for seg in segments:
                start_s = int(seg.get("start", 0))
                end_s = int(seg.get("end", 0))
                start_m, start_sec = divmod(start_s, 60)
                end_m, end_sec = divmod(end_s, 60)
                time_str = f"{start_m:02d}:{start_sec:02d} - {end_m:02d}:{end_sec:02d}"
                seg_text = seg.get("text", "").strip()
                if seg_text:
                    lines.append(f"- `[{time_str}]` {seg_text}")
            res.text = "\n".join(lines)
            res.status = "ocr"
            res.ocr_engine = "whisper-tiny"
            res.language = lang
            res.notes.append("transcribed via Whisper ASR with timestamp breakdown")
            return res
        elif transcription.get("text", "").strip():
            res.text = f"# Media Transcript: {display_name}\n\n" + transcription["text"].strip()
            res.status = "ocr"
            res.ocr_engine = "whisper-tiny"
            res.notes.append("transcribed via Whisper ASR")
            return res
    except Exception:
        pass

    # 2. Try ffmpeg / ffprobe metadata / track extraction if available
    try:
        import json
        import shutil
        import subprocess

        if shutil.which("ffprobe"):
            proc = subprocess.run(
                [
                    "ffprobe",
                    "-v",
                    "quiet",
                    "-print_format",
                    "json",
                    "-show_format",
                    "-show_streams",
                    str(path),
                ],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if proc.returncode == 0:
                info = json.loads(proc.stdout)
                fmt = info.get("format", {})
                tags = fmt.get("tags", {})
                duration = float(fmt.get("duration", 0.0))
                dur_m, dur_s = divmod(int(duration), 60)
                streams = info.get("streams", [])
                lines = [
                    f"# Media File: {display_name}",
                    f"- **Duration**: {dur_m:02d}:{dur_s:02d} ({duration:.2f}s)",
                    f"- **Format**: {fmt.get('format_long_name', ext[1:].upper())}",
                ]
                if fmt.get("bit_rate"):
                    lines.append(f"- **Bitrate**: {int(fmt['bit_rate']) // 1000} kbps")
                if tags.get("title"):
                    lines.append(f"- **Title**: {tags['title']}")
                if tags.get("artist"):
                    lines.append(f"- **Artist**: {tags['artist']}")
                lines.append("\n## Streams Breakdown")
                for s in streams:
                    codec_type = s.get("codec_type")
                    codec_name = s.get("codec_name")
                    if codec_type and codec_name:
                        details = []
                        if s.get("width") and s.get("height"):
                            details.append(f"{s['width']}x{s['height']}")
                        if s.get("r_frame_rate"):
                            details.append(f"{s['r_frame_rate']} fps")
                        if s.get("sample_rate"):
                            details.append(f"{s['sample_rate']} Hz")
                        if s.get("channels"):
                            details.append(f"{s['channels']} ch")
                        extra = f" ({', '.join(details)})" if details else ""
                        lines.append(f"- Stream ({codec_type.capitalize()}): `{codec_name}`{extra}")
                res.text = "\n".join(lines)
                res.status = "native"
                res.notes.append("extracted media stream metadata via ffprobe")
                return res
    except Exception:
        pass

    # 3. Clean fallback metadata representation
    fmt_name = ext.lstrip(".").upper()
    media_kind = "Audio" if ext in AUDIO_EXTS else "Video"
    res.text = f"# Media Resource: {display_name}\nFormat: {fmt_name}\nType: {media_kind} recording\nStatus: Media file registered and ready for playback."
    res.status = "native"
    res.notes.append("media file registered")
    return res


# ── Public API ─────────────────────────────────────────────────────────────────


def extract_text(
    path: Path,
    content_type: str = "",
    filename: str = "",
    document_id: int | None = None,
) -> OcrResult:
    """Route a stored file to the right extractor and return text + status.

    Tier 0 (preferred): Docling — structure-aware, table-preserving parser with ASR.
      - PDFs exceeding settings.pdf_batch_size pages are split into batches
        and processed in parallel using a persistent subprocess pool.
    Fallback: legacy Tesseract / pypdf / office / media parsers.

    ``document_id`` must be supplied when calling from the ingestion worker so
    the parallel path can check the cancellation registry between batch polls.
    """
    ext = Path(filename).suffix.lower() if filename else path.suffix.lower()

    # Cheap preflight budgets protect optional parsers and OCR from hostile
    # page/pixel/text expansion.  The file has already passed signature checks.
    if ext == ".pdf" and _HAS_PYPDF:
        try:
            if len(PdfReader(str(path)).pages) > settings.max_document_pages:
                return _finalize_result(
                    OcrResult(
                        status="skipped",
                        extractor_name="preflight",
                        notes=["page budget exceeded"],
                    )
                )
        except Exception:
            pass
    if ext in IMAGE_EXTS and _HAS_TESSERACT_LIB:
        try:
            with Image.open(path) as image:
                if image.width * image.height > settings.max_image_pixels:
                    return _finalize_result(
                        OcrResult(
                            status="skipped",
                            page_count=1,
                            extractor_name="preflight",
                            notes=["pixel budget exceeded"],
                        )
                    )
        except Exception:
            pass

    if ext == ".vtt":
        result = _extract_vtt(path, filename=filename)
        return _finalize_result(result)

    if ext in MEDIA_EXTS:
        result = _extract_media_fallback(path, ext, filename=filename)
        return _finalize_result(result)

    # Docling handles PDF, DOCX, PPTX, XLSX, HTML, Markdown, CSV, and image
    # formats in one unified pipeline.
    if _docling_available() and ext in DOCLING_EXTS:
        # ── PDF: route large documents through the parallel batch path ──────
        if ext == ".pdf" and _HAS_PYPDF and document_id is not None:
            try:
                page_count = len(PdfReader(str(path)).pages)
            except Exception:
                page_count = 0
            if page_count > settings.pdf_batch_size:
                try:
                    result = _extract_with_docling_parallel(path, document_id)
                    if result.text and result.status != "error":
                        if len(result.text) > settings.max_extracted_text_chars:
                            result.text = result.text[: settings.max_extracted_text_chars]
                            result.notes.append("text budget exceeded; output truncated")
                            result.quality_signals["text_truncated"] = True
                        return _finalize_result(result)
                    result.notes.append("parallel batch extraction empty/error; falling back")
                except (ExtractionCancelled, ExtractionBatchFailed):
                    raise  # propagate cancellation and exhausted-retry failures
                except Exception as _pe:
                    logging.getLogger(__name__).warning(
                        "Parallel batch extraction failed for document %s: %s — falling back to single-pass",
                        document_id, _pe,
                    )

        # ── Single-pass Docling (non-PDF or small PDF or no document_id) ────
        result = _extract_with_docling(path)
        # Only fall back if Docling actually failed (no text extracted + error).
        if result.text and result.status != "error":
            return _finalize_result(result)
        # Docling failed — record the note and fall through to the fallback path.
        result.notes.append("falling back to legacy/media extractor")

    # Legacy / media fallback path.
    if ext in MEDIA_EXTS:
        result = _extract_media_fallback(path, ext, filename=filename)
    elif ext == ".pdf":
        result = _extract_pdf_legacy(path)
    elif ext in IMAGE_EXTS:
        result = _extract_image_legacy(path)
    else:
        result = _extract_office_or_text(path, ext)
    if len(result.text) > settings.max_extracted_text_chars:
        result.text = result.text[: settings.max_extracted_text_chars]
        result.notes.append("text budget exceeded; output truncated")
        result.quality_signals["text_truncated"] = True
    return _finalize_result(result)


def engine_status() -> dict:
    effective_languages, missing_languages = (
        _effective_tesseract_languages() if _HAS_TESSERACT_LIB else ([], list(settings.ocr_languages))
    )
    return {
        "docling": _docling_available(),
        "tesseract": _tesseract_available(),
        "pymupdf": _HAS_FITZ,
        "pypdf": _HAS_PYPDF,
        "ocr_languages": effective_languages,
        "missing_ocr_languages": missing_languages,
    }


def passive_engine_status() -> dict:
    """Report imported capabilities without invoking binaries or loading models."""

    return {
        "docling": _docling_available(),
        "tesseract": _HAS_TESSERACT_LIB,
        "pymupdf": _HAS_FITZ,
        "pypdf": _HAS_PYPDF,
        "requested_ocr_languages": list(settings.ocr_languages),
    }


def warm_docling() -> None:
    """Pre-load the Docling DocumentConverter in a background thread.

    Call this from the FastAPI startup event so the first uploaded file never
    triggers the expensive 770-weight cold-start. The converter is cached as a
    module-level singleton and reused for every subsequent request.
    """
    import threading

    context = worker_context("docling-warm")

    def _load() -> None:
        with bound_request_context(context):
            if not _docling_available():
                emit_event(
                    "worker.model_warm.completed",
                    context=context,
                    component="extraction",
                    operation="docling",
                    outcome="disabled",
                )
                return
            emit_event(
                "worker.model_warm.started",
                context=context,
                component="extraction",
                operation="docling",
            )
            conv = _get_docling_converter()
            emit_event(
                "worker.model_warm.completed",
                context=context,
                component="extraction",
                operation="docling",
                outcome="success" if conv is not None else "unavailable",
            )

    t = threading.Thread(target=_load, name="docling-warm", daemon=True)
    try:
        t.start()
    except Exception as exc:
        emit_event(
            "worker.model_warm.rejected",
            level=logging.ERROR,
            context=context,
            component="extraction",
            operation="thread_start",
            outcome="error",
            error=exc,
        )
        raise
